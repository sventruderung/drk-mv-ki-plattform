"""Textextraktion aus Upload-Formaten (§3.2).

Unterstützt: PDF, Word (doc/docx/docm/rtf/odt), Excel (xls/xlsx/xlsm/ods/csv),
PowerPoint (ppt/pptx), E-Mail (eml/msg), Text (txt/md), HTML.

Liefert Liste von (seite, text)-Tupeln; Seite=None wenn das Format keine
Seiten kennt. PDF nutzt Seitenzahlen, PowerPoint die Foliennummer.
"""

import io
import subprocess
import tempfile
from html.parser import HTMLParser

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pypdf import PdfReader

SUPPORTED_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.ms-word.document.macroEnabled.12": "docx",   # .docm
    "application/msword": "doc",
    "application/rtf": "rtf",
    "text/rtf": "rtf",
    "application/vnd.oasis.opendocument.text": "odt",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel.sheet.macroEnabled.12": "xlsx",     # .xlsm
    "application/vnd.ms-excel": "xls",
    "application/vnd.oasis.opendocument.spreadsheet": "ods",
    "text/csv": "txt",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "ppt",
    "message/rfc822": "eml",
    "application/vnd.ms-outlook": "msg",
    "text/html": "html",
    "text/markdown": "txt",
    "text/plain": "txt",
    "image/png": "image",
    "image/jpeg": "image",
    "image/tiff": "image",
}

OCR_MAX_PAGES = 300  # Schutz vor versehentlichen Mammut-Scans


def extract_text(data: bytes, content_type: str) -> list[tuple[int | None, str]]:
    kind = SUPPORTED_TYPES.get(content_type)
    if kind is None:
        raise ValueError(
            f"Nicht unterstütztes Format: {content_type}. "
            "Unterstützt: PDF, Word, Excel, PowerPoint, OpenOffice, "
            "RTF, E-Mail (EML/MSG), HTML, CSV, TXT"
        )
    match kind:
        case "pdf":
            return _extract_pdf(data)
        case "docx":
            return _extract_docx(data)
        case "doc":
            return _extract_doc(data)
        case "rtf":
            return _extract_rtf(data)
        case "odt":
            return _extract_odf_text(data)
        case "xlsx":
            return _extract_xlsx(data)
        case "xls":
            return _extract_xls(data)
        case "ods":
            return _extract_ods(data)
        case "pptx":
            return _extract_pptx(data)
        case "ppt":
            return _extract_ppt(data)
        case "eml":
            return _extract_eml(data)
        case "msg":
            return _extract_msg(data)
        case "html":
            return _extract_html(data)
        case "image":
            return _extract_image(data)
        case "txt":
            return [(None, data.decode("utf-8", errors="replace"))]
    raise AssertionError("unreachable")


def _extract_pdf(data: bytes) -> list[tuple[int | None, str]]:
    reader = PdfReader(io.BytesIO(data))
    pages = [
        (i + 1, text)
        for i, page in enumerate(reader.pages)
        if (text := page.extract_text() or "").strip()
    ]
    if not pages and len(reader.pages) > 0:
        # Keine Textebene — vermutlich ein Scan: Texterkennung (OCR)
        return _ocr_pdf(data)
    return pages


def _ocr_pdf(data: bytes) -> list[tuple[int | None, str]]:
    """Gescanntes PDF: Seiten rendern und per Tesseract (deutsch) erkennen."""
    import pypdfium2 as pdfium
    import pytesseract

    pdf = pdfium.PdfDocument(data)
    if len(pdf) > OCR_MAX_PAGES:
        raise ValueError(
            f"Scan mit {len(pdf)} Seiten — Maximum für Texterkennung: "
            f"{OCR_MAX_PAGES}. Bitte aufteilen."
        )
    pages: list[tuple[int | None, str]] = []
    for i in range(len(pdf)):
        image = pdf[i].render(scale=2.0).to_pil()  # ~150 dpi
        text = pytesseract.image_to_string(image, lang="deu").strip()
        if text:
            pages.append((i + 1, text))
    if not pages:
        raise ValueError(
            "Auch per Texterkennung (OCR) kein Text gefunden — "
            "Scanqualität prüfen."
        )
    return pages


def _extract_image(data: bytes) -> list[tuple[int | None, str]]:
    """Bilddatei (PNG/JPG/TIFF): Texterkennung per Tesseract (deutsch)."""
    import pytesseract
    from PIL import Image

    try:
        image = Image.open(io.BytesIO(data))
    except Exception:
        raise ValueError("Bilddatei konnte nicht gelesen werden.")
    text = pytesseract.image_to_string(image, lang="deu").strip()
    if not text:
        raise ValueError(
            "Per Texterkennung (OCR) kein Text im Bild gefunden — "
            "Scanqualität prüfen."
        )
    return [(None, text)]


def _extract_docx(data: bytes) -> list[tuple[int | None, str]]:
    doc = DocxDocument(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    # Tabellen mitnehmen — in Formularen steckt der Inhalt oft nur dort
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts)
    return [(None, text)] if text else []


def _extract_doc(data: bytes) -> list[tuple[int | None, str]]:
    """Alt-Word (.doc) via antiword — liefert reinen Text."""
    with tempfile.NamedTemporaryFile(suffix=".doc") as tmp:
        tmp.write(data)
        tmp.flush()
        result = subprocess.run(
            ["antiword", "-m", "UTF-8.txt", tmp.name],
            capture_output=True, timeout=60,
        )
    if result.returncode != 0:
        raise ValueError(
            "Alt-Word-Datei konnte nicht gelesen werden — "
            "bitte als .docx speichern und erneut hochladen."
        )
    text = result.stdout.decode("utf-8", errors="replace").strip()
    return [(None, text)] if text else []


def _extract_rtf(data: bytes) -> list[tuple[int | None, str]]:
    from striprtf.striprtf import rtf_to_text

    try:
        text = rtf_to_text(data.decode("utf-8", errors="replace")).strip()
    except Exception:
        raise ValueError("RTF-Datei konnte nicht gelesen werden.")
    return [(None, text)] if text else []


def _extract_odf_text(data: bytes) -> list[tuple[int | None, str]]:
    """OpenDocument Text (.odt) via odfpy."""
    from odf import teletype, text as odf_text
    from odf.opendocument import load

    doc = load(io.BytesIO(data))
    parts = [
        teletype.extractText(p)
        for p in doc.getElementsByType(odf_text.P)
        if teletype.extractText(p).strip()
    ]
    text = "\n".join(parts)
    return [(None, text)] if text else []


def _extract_xlsx(data: bytes) -> list[tuple[int | None, str]]:
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        rows = [
            " | ".join(str(c) for c in row if c is not None)
            for row in sheet.iter_rows(values_only=True)
            if any(c is not None for c in row)
        ]
        if rows:
            parts.append(f"[Tabelle: {sheet.title}]\n" + "\n".join(rows))
    text = "\n\n".join(parts)
    return [(None, text)] if text else []


def _extract_xls(data: bytes) -> list[tuple[int | None, str]]:
    """Alt-Excel (.xls) via xlrd."""
    import xlrd

    try:
        wb = xlrd.open_workbook(file_contents=data)
    except xlrd.XLRDError:
        raise ValueError(
            "Alt-Excel-Datei konnte nicht gelesen werden — "
            "bitte als .xlsx speichern und erneut hochladen."
        )
    parts: list[str] = []
    for sheet in wb.sheets():
        rows = [
            " | ".join(str(c) for c in sheet.row_values(r) if c not in ("", None))
            for r in range(sheet.nrows)
            if any(c not in ("", None) for c in sheet.row_values(r))
        ]
        if rows:
            parts.append(f"[Tabelle: {sheet.name}]\n" + "\n".join(rows))
    text = "\n\n".join(parts)
    return [(None, text)] if text else []


def _extract_ods(data: bytes) -> list[tuple[int | None, str]]:
    """OpenDocument Spreadsheet (.ods) via odfpy."""
    from odf import teletype
    from odf.opendocument import load
    from odf.table import Table, TableCell, TableRow

    doc = load(io.BytesIO(data))
    parts: list[str] = []
    for table in doc.getElementsByType(Table):
        rows = []
        for row in table.getElementsByType(TableRow):
            cells = [
                teletype.extractText(c).strip()
                for c in row.getElementsByType(TableCell)
            ]
            cells = [c for c in cells if c]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            name = table.getAttribute("name") or "Tabelle"
            parts.append(f"[Tabelle: {name}]\n" + "\n".join(rows))
    text = "\n\n".join(parts)
    return [(None, text)] if text else []


def _extract_pptx(data: bytes) -> list[tuple[int | None, str]]:
    """PowerPoint (.pptx) — Foliennummer dient als Seitenangabe für Zitate."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    pages: list[tuple[int | None, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            texts.append(f"[Notizen] {slide.notes_slide.notes_text_frame.text}")
        if texts:
            pages.append((i, "\n".join(texts)))
    return pages


def _extract_ppt(data: bytes) -> list[tuple[int | None, str]]:
    """Alt-PowerPoint (.ppt) via catppt (Paket catdoc)."""
    with tempfile.NamedTemporaryFile(suffix=".ppt") as tmp:
        tmp.write(data)
        tmp.flush()
        result = subprocess.run(
            ["catppt", "-dutf-8", tmp.name], capture_output=True, timeout=60
        )
    if result.returncode != 0:
        raise ValueError(
            "Alt-PowerPoint-Datei konnte nicht gelesen werden — "
            "bitte als .pptx speichern und erneut hochladen."
        )
    text = result.stdout.decode("utf-8", errors="replace").strip()
    return [(None, text)] if text else []


def _extract_eml(data: bytes) -> list[tuple[int | None, str]]:
    """E-Mail (.eml) — Kopfzeilen + Textkörper."""
    import email
    from email import policy

    m = email.message_from_bytes(data, policy=policy.default)
    header = "\n".join(
        f"{k}: {m.get(k)}" for k in ("From", "To", "Date", "Subject") if m.get(k)
    )
    body = ""
    part = m.get_body(preferencelist=("plain", "html"))
    if part is not None:
        content = part.get_content()
        body = _strip_html(content) if part.get_content_type() == "text/html" else content
    text = f"{header}\n\n{body}".strip()
    return [(None, text)] if text else []


def _extract_msg(data: bytes) -> list[tuple[int | None, str]]:
    """Outlook-E-Mail (.msg) via extract_msg."""
    import extract_msg

    with tempfile.NamedTemporaryFile(suffix=".msg") as tmp:
        tmp.write(data)
        tmp.flush()
        try:
            m = extract_msg.openMsg(tmp.name)
        except Exception:
            raise ValueError("Outlook-Datei (.msg) konnte nicht gelesen werden.")
        header = "\n".join(
            f"{label}: {value}"
            for label, value in (
                ("From", m.sender), ("To", m.to),
                ("Date", m.date), ("Subject", m.subject),
            ) if value
        )
        text = f"{header}\n\n{m.body or ''}".strip()
        m.close()
    return [(None, text)] if text else []


class _HTMLTextParser(HTMLParser):
    SKIP = {"script", "style"}

    def __init__(self):
        super().__init__()
        self.parts: list[str] = []
        self._skip = 0

    def handle_starttag(self, tag, attrs):
        if tag in self.SKIP:
            self._skip += 1

    def handle_endtag(self, tag):
        if tag in self.SKIP and self._skip:
            self._skip -= 1

    def handle_data(self, d):
        if not self._skip and d.strip():
            self.parts.append(d.strip())


def _strip_html(html: str) -> str:
    parser = _HTMLTextParser()
    parser.feed(html)
    return "\n".join(parser.parts)


def _extract_html(data: bytes) -> list[tuple[int | None, str]]:
    text = _strip_html(data.decode("utf-8", errors="replace")).strip()
    return [(None, text)] if text else []
